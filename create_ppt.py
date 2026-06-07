import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_modern_slide(prs, title_text=""):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Add dark background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(15, 23, 42)  # Slate 900
    bg.line.fill.background()
    
    # Add title if provided
    if title_text:
        tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1))
        tf = tx_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(248, 250, 252)  # Slate 50
    
    return slide

def add_body_text(slide, text, left, top, width, height, font_size=20, align=PP_ALIGN.LEFT):
    tx_box = slide.shapes.add_textbox(left, top, width, height)
    tf = tx_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = RGBColor(203, 213, 225)  # Slate 300
    p.alignment = align
    return tx_box

def main():
    prs = Presentation()
    
    # Page 1: Title & Team
    slide1 = create_modern_slide(prs)
    
    # VINGEL Title
    tx_box = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    tf = tx_box.text_frame
    p = tf.paragraphs[0]
    p.text = "VINGEL"
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = RGBColor(129, 140, 248)  # Indigo 400
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    add_body_text(slide1, "Market Simulator & On-Chain Analytics", Inches(1), Inches(3.2), Inches(8), Inches(0.5), 24, PP_ALIGN.CENTER)
    
    # Team Members
    team_text = "Team Members:\n1) Sumit kumar\n2) Aman Mishra\n3) Nakul Varshney"
    add_body_text(slide1, team_text, Inches(3.5), Inches(4.5), Inches(4), Inches(2), 22, PP_ALIGN.LEFT)

    # Page 2: AI Market Segmentation
    slide2 = create_modern_slide(prs, "AI-Powered Market Segmentation")
    add_body_text(slide2, "• Vingel uses advanced AI to generate distinct customer archetypes.\n• Captures dynamic behavioral parameters like price sensitivity and tech affinity.\n• Replaces static guesses with realistic data points.", Inches(0.5), Inches(1.5), Inches(4.5), Inches(5), 22)
    img_path1 = r"C:\Users\Nakul\.gemini\antigravity-ide\brain\0c7c17cd-88a0-4f1a-aaa5-ecbc6c1b8c59\modern_ai_segmentation_1780831110658.png"
    if os.path.exists(img_path1):
        slide2.shapes.add_picture(img_path1, Inches(5), Inches(1.5), height=Inches(5))

    # Page 3: Neo4j Population Graph
    slide3 = create_modern_slide(prs, "Real-Time Population Visualization")
    add_body_text(slide3, "• Synthetic populations modeled as connected network graphs.\n• High-performance interactive clustering with vibrant visuals.\n• Visualizes real-time conversion and churn events.", Inches(0.5), Inches(1.5), Inches(4.5), Inches(5), 22)
    img_path2 = r"C:\Users\Nakul\.gemini\antigravity-ide\brain\0c7c17cd-88a0-4f1a-aaa5-ecbc6c1b8c59\modern_network_graph_1780831123386.png"
    if os.path.exists(img_path2):
        slide3.shapes.add_picture(img_path2, Inches(5), Inches(1.5), height=Inches(5))

    # Page 4: Monad Blockchain Anchoring
    slide4 = create_modern_slide(prs, "On-Chain Blockchain Anchoring")
    add_body_text(slide4, "• Immutable proof-of-simulation stored on the Monad testnet.\n• End-to-end cryptographic integrity via SHA-256 block chains.\n• Secure, decentralized, and gated by wallet balances.", Inches(0.5), Inches(1.5), Inches(4.5), Inches(5), 22)
    img_path3 = r"C:\Users\Nakul\.gemini\antigravity-ide\brain\0c7c17cd-88a0-4f1a-aaa5-ecbc6c1b8c59\modern_blockchain_ledger_1780831109303.png"
    if os.path.exists(img_path3):
        slide4.shapes.add_picture(img_path3, Inches(5), Inches(1.5), height=Inches(5))

    # Page 5: Thank You
    slide5 = create_modern_slide(prs)
    tx_box5 = slide5.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
    tf5 = tx_box5.text_frame
    p5 = tf5.paragraphs[0]
    p5.text = "Thank You!"
    p5.font.size = Pt(64)
    p5.font.bold = True
    p5.font.color.rgb = RGBColor(129, 140, 248)
    p5.alignment = PP_ALIGN.CENTER

    prs.save("Vingel_Presentation.pptx")
    print("Vingel_Presentation.pptx generated successfully.")

if __name__ == "__main__":
    main()
